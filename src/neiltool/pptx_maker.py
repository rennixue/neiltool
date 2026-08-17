import asyncio
import logging
import re
from collections.abc import Sequence
from io import BytesIO
from typing import TypedDict

from lxml import etree  # type: ignore
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE
from pptx.shapes.autoshape import Shape
from pptx.slide import Slide
from pptx.util import Inches, Length, Pt

from .models.operation import SlideData

logger = logging.getLogger(__name__)


def make_pptx(data: SlideData) -> bytes | None:
    try:
        bytes_ = PresentationTemplate().render(data)
        return bytes_
    except Exception as exc:
        logger.error("fail to make pptx: %r", exc)
    return None


async def make_pptx_async(data: SlideData) -> bytes | None:
    return await asyncio.to_thread(make_pptx, data)


class LtwhDict(TypedDict):
    left: Length
    top: Length
    width: Length
    height: Length


def ltwh(left: float, top: float, width: float, height: float) -> LtwhDict:
    return {
        "left": Inches(left),
        "top": Inches(top),
        "width": Inches(width),
        "height": Inches(height),
    }


def edit_rounded_corner(mut_shape: Shape, value: int) -> None:
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    prst_geom = mut_shape._element.find(".//a:prstGeom", nsmap)  # type: ignore
    if prst_geom is None:
        return
    av_lst = prst_geom.find("a:avLst", nsmap)  # type: ignore
    if av_lst is None:
        av_lst = etree.SubElement(prst_geom, "{http://schemas.openxmlformats.org/drawingml/2006/main}avLst")  # type: ignore
    else:
        av_lst.clear()  # type: ignore
    gd = etree.SubElement(av_lst, "{http://schemas.openxmlformats.org/drawingml/2006/main}gd")  # type: ignore
    gd.set("name", "adj")  # type: ignore
    gd.set("fmla", f"val {value}")  # type: ignore


class PresentationTemplate:
    COLORS = {
        "text": RGBColor(0, 0, 0),
        "emph": RGBColor(0, 64, 128),
        "mute": RGBColor(128, 128, 128),
    }
    FONT_SIZES = {
        "small": Pt(10),
        "normal": Pt(11),
        "large": Pt(12),
    }

    def render(self, data: SlideData) -> bytes:
        presentation = Presentation()
        presentation.slide_width = Inches(40 / 3)
        presentation.slide_height = Inches(7.5)
        slide_layout = presentation.slide_layouts[6]  # blank
        slide = presentation.slides.add_slide(slide_layout)
        self.render_slide(slide, data)
        bytes_io = BytesIO()
        presentation.save(bytes_io)
        return bytes_io.getvalue()

    def render_slide(self, mut_slide: Slide, data: SlideData) -> None:
        self.render_title(mut_slide, data)
        self.render_phases(mut_slide, data)
        self.render_body(mut_slide, data)
        self.render_footer(mut_slide, data)

    def add_textbox(self, mut_slide: Slide, ltwh_dict: LtwhDict, texts: Sequence[tuple[str, Sequence[str]]]) -> Shape:
        shape = mut_slide.shapes.add_textbox(**ltwh_dict)
        text_frame = shape.text_frame
        for i, (heading, paras) in enumerate(texts):
            if i == 0:
                if heading:
                    paragraph = text_frame.paragraphs[0]
                    paragraph.text = heading
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = self.COLORS["emph"]
                    paragraph.font.size = self.FONT_SIZES["large"]
                    paragraph.space_after = Pt(6)
                for _ in range(len(paras) - (0 if heading else 1)):
                    text_frame.add_paragraph()
                for paragraph, s in zip(text_frame.paragraphs[(1 if heading else 0) :], paras):
                    paragraph.text = s
                    paragraph.font.size = self.FONT_SIZES["normal"]
            else:
                if heading:
                    paragraph = text_frame.add_paragraph()
                    paragraph.text = heading
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = self.COLORS["emph"]
                    paragraph.font.size = self.FONT_SIZES["large"]
                    paragraph.space_before = Pt(12)
                    paragraph.space_after = Pt(6)
                for s in paras:
                    paragraph = text_frame.add_paragraph()
                    paragraph.text = s
                    paragraph.font.size = self.FONT_SIZES["normal"]
        return shape

    def add_rect(
        self,
        mut_slide: Slide,
        ltwh_dict: LtwhDict,
        border_color: RGBColor | None = None,
        bg_color: RGBColor | None = None,
    ) -> Shape:
        shape = mut_slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, **ltwh_dict)
        shape.shadow.inherit = False
        shape.line.fill.solid()
        shape.line.width = Pt(1)
        if border_color:
            shape.line.color.rgb = border_color
        else:
            shape.line.color.rgb = RGBColor(224, 224, 224)
        if bg_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = bg_color
        else:
            shape.fill.background()
        return shape

    def render_title(self, mut_slide: Slide, data: SlideData) -> None:
        shape = self.add_textbox(
            mut_slide,
            ltwh(0.5, 0.3, 6.0, 0.8),
            [("上课方案与总结", ["基于当前课程信息、订单需求和破冰课内容整理出的后续上课路径"])],
        )
        text_frame = shape.text_frame
        paragraph = text_frame.paragraphs[0]
        paragraph.font.color.rgb = self.COLORS["text"]
        paragraph.font.size = Pt(24)
        paragraph = text_frame.paragraphs[1]
        paragraph.font.size = Pt(12)

    def render_phases(self, mut_slide: Slide, data: SlideData) -> None:
        labels = data.phases
        if not labels:
            return
        radius = 0.2
        interval = 1.2
        start = 6.7 - (len(labels) / 2) * interval - radius
        vertical_center = 1.4
        line = mut_slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT,
            Inches(start + radius),
            Inches(vertical_center),
            Inches(start + radius + interval * (len(labels) - 1)),
            Inches(vertical_center),
        )
        line.line.color.rgb = RGBColor(224, 224, 224)
        line.line.width = Pt(2)
        line.shadow.inherit = False
        for i, label in enumerate(labels):
            shape = mut_slide.shapes.add_shape(
                MSO_SHAPE.OVAL, **ltwh(start + interval * i, vertical_center - radius, radius * 2, radius * 2)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.COLORS["emph"]
            shape.line.fill.background()
            shape.shadow.inherit = False
            text_frame = shape.text_frame
            paragraph = text_frame.paragraphs[0]
            paragraph.text = str(i + 1)
            paragraph.font.size = Pt(9)
            paragraph.font.color.rgb = RGBColor(255, 255, 255)

            shape = self.add_textbox(
                mut_slide, ltwh(start - 0.18 + interval * i, vertical_center + radius, 1.0, 0.3), [(label, [])]
            )
            text_frame = shape.text_frame
            paragraph = text_frame.paragraphs[0]
            paragraph.font.color.rgb = self.COLORS["text"]
            paragraph.font.size = self.FONT_SIZES["small"]

    def render_footer(self, mut_slide: Slide, data: SlideData) -> None:
        shape = self.add_rect(mut_slide, ltwh(0.5, 7.0, 12.4, 0.3), RGBColor(64, 192, 64), RGBColor(248, 255, 248))
        shape = self.add_textbox(mut_slide, ltwh(0.5, 7.0, 12.4, 0.3), [("", [" "])])
        text_frame = shape.text_frame
        paragraph = text_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = "课后行动\u3000"
        run.font.color.rgb = RGBColor(64, 192, 64)
        run.font.bold = True
        run = paragraph.add_run()
        run.text = data.afterclass

    def render_body(self, mut_slide: Slide, data: SlideData) -> None:
        ltwh_left = ltwh(0.5, 2.0, 4.0, 4.8)
        ltwh_mid = LtwhDict(**{**ltwh_left, **{"left": Inches(4.7)}})  # type: ignore
        ltwh_right = LtwhDict(**{**ltwh_left, **{"left": Inches(8.9)}})  # type: ignore

        shape = self.add_rect(mut_slide, ltwh(0.5, 2.0, 4.0, 4.8), None, RGBColor(248, 248, 255))
        edit_rounded_corner(shape, 2000)
        contents: list[str] = []
        if s := data.course_name:
            contents.append(f"课程名称：{s}")
        if s := data.course_code:
            contents.append(f"课程代码：{s}")
        if s := data.assessment:
            contents.append(f"考核项：{s}")
        shape = self.add_textbox(
            mut_slide,
            ltwh_left,
            [
                ("基本信息", []),
                *[("", [it]) for it in contents],
                ("学习目标", [data.what]),
                ("如何达成", [data.how]),
                ("如何安排", [data.why]),
            ],
        )
        shape.text_frame.word_wrap = True

        shape = self.add_rect(mut_slide, ltwh_mid)
        edit_rounded_corner(shape, 2000)
        contents: list[str] = []
        for i, (group, keypoints) in enumerate(data.knowledges, 1):
            contents.append(f"{i}. {group}")
            contents.append(", ".join(keypoints))
        shape = self.add_textbox(mut_slide, ltwh_mid, [("知识点梳理", []), *[("", [it]) for it in contents]])
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        for i, paragraph in enumerate(text_frame.paragraphs[1:]):
            if i % 2 == 0:
                paragraph.font.bold = True
                paragraph.space_before = Pt(6)
            else:
                paragraph.font.size = self.FONT_SIZES["small"]
                paragraph.font.color.rgb = self.COLORS["mute"]

        shape = self.add_rect(mut_slide, ltwh_right, None, RGBColor(248, 248, 255))
        edit_rounded_corner(shape, 2000)
        shape = self.add_textbox(mut_slide, ltwh_right, [("课程将完成", [])])
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        for objective in data.objectives:
            paragraph = text_frame.add_paragraph()
            if m := re.match(r"\d+", objective):
                head, tail = objective[: m.end()], objective[m.end() :]
                run = paragraph.add_run()
                run.text = head + " "
                run.font.size = Pt(16)
                run.font.color.rgb = self.COLORS["emph"]
                run = paragraph.add_run()
                run.text = tail
                paragraph.font.size = self.FONT_SIZES["normal"]
            else:
                paragraph.text = objective
                paragraph.font.size = self.FONT_SIZES["normal"]
        text_frame.add_paragraph()
        paragraph = text_frame.add_paragraph()
        paragraph.text = "后续每节课资料与错题将持续沉淀为考前重点复习文档，具体课时与课次请和老师确认"
        paragraph.font.size = self.FONT_SIZES["normal"]
